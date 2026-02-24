#!/usr/bin/env python3
"""
Test script to verify slide reordering functionality.
"""

from pptx import Presentation
from services.pptx_service import PowerPointService

def test_slide_reordering():
    """Test that slides can be reordered correctly."""
    
    print("=" * 70)
    print("TESTING SLIDE REORDERING")
    print("=" * 70)
    
    # Create a test presentation
    print("\n1. Creating test presentation with 5 slides...")
    prs = Presentation()
    
    for i in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        title = slide.shapes.title
        title.text = f"Slide {i+1}: Test Slide"
    
    print(f"   Created {len(prs.slides)} slides")
    
    # Print original order
    print("\n2. Original slide order:")
    for idx, slide in enumerate(prs.slides):
        if slide.shapes.title:
            print(f"   Position {idx}: {slide.shapes.title.text}")
    
    # Create mock kept_slides list (reorder: 4, 2, 0, 1, 3)
    kept_slides = [
        {'slide': {'index': 4, 'title': 'Slide 5'}, 'score': 10},
        {'slide': {'index': 2, 'title': 'Slide 3'}, 'score': 9},
        {'slide': {'index': 0, 'title': 'Slide 1'}, 'score': 8},
        {'slide': {'index': 1, 'title': 'Slide 2'}, 'score': 7},
        {'slide': {'index': 3, 'title': 'Slide 4'}, 'score': 6},
    ]
    
    # Reorder slides
    print("\n3. Reordering slides by relevance score...")
    service = PowerPointService()
    prs = service.reorder_slides(prs, kept_slides)
    
    # Print new order
    print("\n4. New slide order:")
    expected_order = ["Slide 5", "Slide 3", "Slide 1", "Slide 2", "Slide 4"]
    success = True
    
    for idx, slide in enumerate(prs.slides):
        if slide.shapes.title:
            actual_title = slide.shapes.title.text
            expected_title = f"{expected_order[idx]}: Test Slide"
            status = "✅" if expected_title in actual_title else "❌"
            print(f"   Position {idx}: {actual_title} {status}")
            if expected_title not in actual_title:
                success = False
    
    # Result
    print("\n" + "=" * 70)
    if success:
        print("✅ SLIDE REORDERING TEST PASSED!")
        print("   Slides were successfully reordered by relevance score")
    else:
        print("❌ SLIDE REORDERING TEST FAILED!")
        print("   Slides were not in the expected order")
    print("=" * 70)
    
    return success

if __name__ == "__main__":
    test_slide_reordering()

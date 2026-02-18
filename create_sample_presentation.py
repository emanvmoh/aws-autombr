from pptx import Presentation
from pptx.util import Inches, Pt

def create_sample_mbr():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Monthly Business Review"
    subtitle.text = "Q1 2026 - Customer Name"
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Executive Summary"
    body = slide.placeholders[1].text_frame
    body.text = "Key highlights from this quarter:\n• Infrastructure overview\n• Cost analysis\n• Support summary\n• Upcoming initiatives"
    
    # Slide 3: AWS Cost Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AWS Cost Overview"
    body = slide.placeholders[1].text_frame
    body.text = "Total spend and breakdown by service\n• EC2 instances\n• RDS databases\n• S3 storage\n• Data transfer"
    
    # Slide 4: Cost Optimization Opportunities
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Cost Optimization Opportunities"
    body = slide.placeholders[1].text_frame
    body.text = "Identified savings opportunities:\n• Reserved Instance recommendations\n• Right-sizing opportunities\n• S3 lifecycle policies\n• Unused resources"
    
    # Slide 5: Support Case Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Support Case Summary"
    body = slide.placeholders[1].text_frame
    body.text = "Recent support activity:\n• Open cases: X\n• Resolved cases: Y\n• Average resolution time\n• Key issues addressed"
    
    # Slide 6: Well-Architected Review
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Well-Architected Review"
    body = slide.placeholders[1].text_frame
    body.text = "Framework pillars assessment:\n• Operational Excellence\n• Security\n• Reliability\n• Performance Efficiency\n• Cost Optimization"
    
    # Slide 7: Security & Compliance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Security & Compliance"
    body = slide.placeholders[1].text_frame
    body.text = "Security posture:\n• IAM best practices\n• Encryption status\n• Compliance requirements\n• Security Hub findings"
    
    # Slide 8: Innovation Opportunities
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Innovation Opportunities"
    body = slide.placeholders[1].text_frame
    body.text = "New AWS services to consider:\n• Serverless technologies\n• Machine learning services\n• Container orchestration\n• Analytics platforms"
    
    # Slide 9: Roadmap & Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Roadmap & Next Steps"
    body = slide.placeholders[1].text_frame
    body.text = "Upcoming initiatives:\n• Q2 priorities\n• Technical deep dives\n• Training sessions\n• Architecture reviews"
    
    # Slide 10: Questions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Questions?"
    
    prs.save('sample_mbr_presentation.pptx')
    print("Sample MBR presentation created: sample_mbr_presentation.pptx")

if __name__ == '__main__':
    create_sample_mbr()

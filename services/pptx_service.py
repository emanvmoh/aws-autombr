from pptx import Presentation
from datetime import datetime

class PowerPointService:
    def load_presentation(self, filepath):
        return Presentation(filepath)
    
    def extract_slide_content(self, prs):
        slides_data = []
        for idx, slide in enumerate(prs.slides):
            slide_data = {'index': idx, 'title': '', 'content': [], 'notes': ''}
            if slide.shapes.title:
                slide_data['title'] = slide.shapes.title.text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    slide_data['content'].append(shape.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_data['notes'] = slide.notes_slide.notes_text_frame.text
            slides_data.append(slide_data)
        return slides_data
    
    def reorder_slides(self, prs, kept_slides):
        """
        Reorder slides by moving them within the presentation.
        
        Args:
            prs: Original presentation
            kept_slides: List of slide items with 'slide' dict and 'score', sorted by relevance
            
        Returns:
            Modified presentation with reordered slides
        """
        # Get the XML element that contains all slides
        sldIdLst = prs.slides._sldIdLst
        
        # Create mapping of original indices to slide IDs
        slide_ids = {i: sldIdLst[i] for i in range(len(prs.slides))}
        
        # Build new order: kept slides sorted by score, then removed slides at end
        new_order_indices = [item['slide']['index'] for item in kept_slides]
        
        # Reorder by removing all and re-adding in new order
        for slide_id in list(sldIdLst):
            sldIdLst.remove(slide_id)
        
        for idx in new_order_indices:
            if idx in slide_ids:
                sldIdLst.append(slide_ids[idx])
        
        return prs
    
    def add_talking_points(self, slide, talking_points):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        if text_frame.text.strip():
            text_frame.text += "\n\n--- TAM TALKING POINTS ---\n"
        else:
            text_frame.text = "--- TAM TALKING POINTS ---\n"
        text_frame.text += talking_points
    
    def save_presentation(self, prs, output_path):
        prs.save(output_path)
        return output_path
    
    def create_change_summary(self, changes):
        summary = "# MBR Presentation Changes Summary\n\n"
        summary += f"Generated: {changes.get('timestamp', datetime.now().isoformat())}\n\n"
        
        if changes.get('removed_slides'):
            summary += "## Slides Removed\n"
            for slide in changes['removed_slides']:
                summary += f"- Slide {slide['index']}: {slide['title']} - {slide['reason']}\n"
            summary += "\n"
        
        if changes.get('reordered'):
            summary += "## Slides Reordered\n"
            for idx, slide in enumerate(changes['reordered'], 1):
                summary += f"{idx}. {slide['title']} (was position {slide['original_index'] + 1})\n"
            summary += "\n"
        
        if changes.get('talking_points_added'):
            summary += f"## Talking Points Added\n{len(changes['talking_points_added'])} slides enhanced\n\n"
        
        if changes.get('customer_context'):
            summary += "## Customer Context\n" + changes['customer_context'] + "\n\n"
        
        return summary
